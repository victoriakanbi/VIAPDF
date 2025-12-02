import streamlit as st
import io
import os
import zipfile
import uuid
import math
import itertools
import string
import time
import concurrent.futures
import xml.etree.ElementTree as ET
from PIL import Image, ImageDraw, ImageFont
from pypdf import PdfReader, PdfWriter
import img2pdf
from pdf2image import convert_from_bytes
from pdf2docx import Converter
import pdfplumber
import pandas as pd
import pikepdf
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.utils import simpleSplit, ImageReader
from reportlab.lib import colors
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

# --- OPTIONAL IMPORTS ---
try:
    from docx2pdf import convert as convert_docx
    HAS_DOCX_SUPPORT = True
except ImportError:
    HAS_DOCX_SUPPORT = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX_SUPPORT = True
except ImportError:
    HAS_PPTX_SUPPORT = False

try:
    import pytesseract
    HAS_OCR_SUPPORT = True
except ImportError:
    HAS_OCR_SUPPORT = False

try:
    from streamlit_drawable_canvas import st_canvas
    HAS_CANVAS_SUPPORT = True
except ImportError:
    HAS_CANVAS_SUPPORT = False

# --- HELPER: AUTO-DETECT EXTERNAL TOOLS ---
def get_local_poppler_path():
    current_dir = os.getcwd()
    for root, dirs, files in os.walk(current_dir):
        if "pdftoppm.exe" in files:
            return root
    return None

def get_local_tesseract_path():
    # 0. User specific path (Priority)
    user_path = r"C:\Users\Surface\OneDrive\Desktop\Add-in project\tesseract.exe"
    if os.path.exists(user_path):
        return user_path

    # 1. Check current directory
    current_dir = os.getcwd()
    for root, dirs, files in os.walk(current_dir):
        if "tesseract.exe" in files:
            return os.path.join(root, "tesseract.exe")

    # 2. Check common Windows Install Paths
    common_paths = [
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
        os.path.join(os.getenv('LOCALAPPDATA', ''), r"Tesseract-OCR\tesseract.exe")
    ]
    
    for path in common_paths:
        if os.path.exists(path):
            return path
    return None

# --- HELPER: GENERATE THUMBNAIL ---
@st.cache_data(show_spinner=False)
def get_page_thumbnail(page_bytes, poppler_path=None, width=200):
    try:
        dpi = 72 if width <= 200 else 150
        common_args = {"first_page": 1, "last_page": 1, "dpi": dpi, "size": (width, None)}
        if poppler_path:
            images = convert_from_bytes(page_bytes, poppler_path=poppler_path, **common_args)
        else:
            images = convert_from_bytes(page_bytes, **common_args)
        return images[0] if images else None
    except Exception:
        return None

# --- HELPER: FONT SELECTOR COMPONENT ---
def font_selector_component(key_prefix):
    font_options = ["Helvetica", "Helvetica-Bold", "Times-Roman", "Times-Bold", "Courier", "Courier-Bold", "Custom (.ttf)"]
    selected_font = st.selectbox("Font Family", font_options, key=f"{key_prefix}_select")
    
    final_font = selected_font
    if selected_font == "Custom (.ttf)":
        uploaded_font = st.file_uploader("Upload .ttf Font", type="ttf", key=f"{key_prefix}_upload")
        if uploaded_font:
            try:
                font_name = f"CustomFont_{key_prefix}"
                temp_path = f"temp_{font_name}.ttf"
                with open(temp_path, "wb") as f:
                    f.write(uploaded_font.read())
                pdfmetrics.registerFont(TTFont(font_name, temp_path))
                final_font = font_name
                st.caption(f"✅ Loaded: {uploaded_font.name}")
            except Exception as e:
                st.error(f"Font Error: {e}")
                final_font = "Helvetica"
    return final_font

# --- HELPER: READ WORDLIST FILES ---
def read_wordlist_file(uploaded_file):
    words = []
    fname = uploaded_file.name.lower()
    try:
        if fname.endswith(".txt"):
            try: content = uploaded_file.getvalue().decode("utf-8")
            except UnicodeDecodeError: content = uploaded_file.getvalue().decode("latin-1")
            words = [line.strip() for line in content.splitlines() if line.strip()]
        elif fname.endswith(".csv"):
            df = pd.read_csv(uploaded_file)
            for col in df.columns: words.extend(df[col].astype(str).str.strip().tolist())
        elif fname.endswith(".xlsx") or fname.endswith(".xls"):
            df = pd.read_excel(uploaded_file)
            for col in df.columns: words.extend(df[col].astype(str).str.strip().tolist())
        elif fname.endswith(".docx"):
            with zipfile.ZipFile(uploaded_file) as zf:
                xml_content = zf.read('word/document.xml')
                tree = ET.fromstring(xml_content)
                for elem in tree.iter():
                    if elem.text and elem.text.strip(): words.append(elem.text.strip())
    except Exception as e: st.error(f"Error reading dictionary file: {e}")
    return list(set(words))

# --- HELPER: TEXT TO PDF ---
def text_to_pdf(text_bytes):
    try: text = text_bytes.decode('utf-8')
    except UnicodeDecodeError: text = text_bytes.decode('latin-1')
    buffer = io.BytesIO()
    c = canvas.Canvas(buffer, pagesize=letter)
    width, height = letter
    text_object = c.beginText(40, height - 40)
    text_object.setFont("Helvetica", 10)
    lines = text.split('\n')
    for line in lines:
        wrapped_lines = simpleSplit(line, "Helvetica", 10, width - 80)
        for wrapped in wrapped_lines:
            if text_object.getY() < 40:
                c.drawText(text_object)
                c.showPage()
                text_object = c.beginText(40, height - 40)
                text_object.setFont("Helvetica", 10)
            text_object.textLine(wrapped)
    c.drawText(text_object)
    c.save()
    buffer.seek(0)
    return buffer.getvalue()

# --- HELPER: PARSE ORDER STRING ---
def parse_order_string(order_str, max_len):
    indices = []
    try:
        parts = [p.strip() for p in order_str.split(',') if p.strip()]
        for p in parts:
            if '-' in p:
                start, end = map(int, p.split('-'))
                indices.extend(range(start - 1, end))
            else:
                indices.append(int(p) - 1)
        return [i for i in indices if 0 <= i < max_len]
    except ValueError:
        return None

# --- HELPER: PDF TO EDITABLE PPTX ---
def create_editable_pptx(pdf_file):
    prs = Presentation()
    with pdfplumber.open(pdf_file) as pdf:
        for page in pdf.pages:
            pdf_w = page.width
            pdf_h = page.height
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            prs.slide_width = Pt(pdf_w)
            prs.slide_height = Pt(pdf_h)
            words = page.extract_words(extra_attrs=["fontname", "size"])
            lines = {}
            for w in words:
                y = round(w['top'], 0) 
                if y not in lines: lines[y] = []
                lines[y].append(w)
            for y in sorted(lines.keys()):
                line_words = lines[y]
                line_text = " ".join([w['text'] for w in line_words])
                x0 = min([w['x0'] for w in line_words])
                top = min([w['top'] for w in line_words])
                width = sum([w['x1'] - w['x0'] for w in line_words]) + (len(line_words) * 3) 
                height = max([w['bottom'] - w['top'] for w in line_words])
                avg_size = sum([w['size'] for w in line_words]) / len(line_words)
                txBox = slide.shapes.add_textbox(Pt(x0), Pt(top), Pt(width), Pt(height))
                tf = txBox.text_frame
                tf.word_wrap = False 
                p = tf.paragraphs[0]
                p.text = line_text
                p.font.size = Pt(avg_size)
                p.font.name = "Arial" 
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()

# --- HELPER: THREADED PASSWORD CHECKER ---
def check_password_batch(filepath, passwords):
    for pw in passwords:
        try:
            with pikepdf.open(filepath, password=str(pw)) as pdf: return str(pw)
        except: continue
    return None

# --- PAGE SETUP ---
st.set_page_config(page_title="VIAPDF", page_icon="🚀", layout="wide", initial_sidebar_state="expanded")

st.markdown("""
<style>
    .main { background-color: #f8f9fa; }
    h1, h2, h3 { color: #002366 !important; font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; }
    .stButton>button { width: 100%; background-color: #002366; color: white; border: none; border-radius: 8px; padding: 0.6rem 1rem; font-weight: 600; transition: all 0.3s ease; box-shadow: 0 2px 5px rgba(0,0,0,0.2); }
    .stButton>button:hover { background-color: #003399; color: white; box-shadow: 0 4px 8px rgba(0,0,0,0.3); transform: translateY(-2px); }
    .stFileUploader { background-color: white; padding: 1rem; border-radius: 10px; border: 2px dashed #002366; }
    .page-card { background: white; padding: 15px; border-radius: 12px; box-shadow: 0 4px 12px rgba(0,0,0,0.08); margin-bottom: 15px; text-align: center; transition: transform 0.2s; border: 1px solid #eef0f2; }
    .page-card:hover { transform: scale(1.02); border-color: #002366; }
    .logo-container { text-align: center; padding: 30px 0; margin-bottom: 20px; background: linear-gradient(180deg, rgba(255,255,255,0) 0%, rgba(0,35,102,0.05) 100%); border-radius: 15px; }
    .logo-via { font-family: 'Arial Black', sans-serif; font-weight: 900; font-size: 72px; color: #002366; letter-spacing: -3px; display: inline-block; }
    .logo-pdf { font-family: 'Arial', sans-serif; font-weight: 300; font-size: 72px; color: #8B0000; letter-spacing: -1px; }
    .logo-sub { font-family: 'Segoe UI', sans-serif; color: #555; font-size: 14px; letter-spacing: 3px; text-transform: uppercase; margin-top: -10px; font-weight: 600; }
    section[data-testid="stSidebar"] { background-color: #ffffff; border-right: 1px solid #e0e0e0; }
</style>
""", unsafe_allow_html=True)

st.markdown("""<div class="logo-container"><span class="logo-via">VIA</span><span class="logo-pdf">PDF</span><div class="logo-sub">Victor's Intelligent Assistant</div></div>""", unsafe_allow_html=True)
st.markdown("<div style='text-align: center; margin-bottom: 30px;'><b>Local. Private. Powerful.</b></div>", unsafe_allow_html=True)

# Session States
if 'page_queue' not in st.session_state: st.session_state['page_queue'] = []
if 'processed_files' not in st.session_state: st.session_state['processed_files'] = set()
if 'extracted_pdf' not in st.session_state: st.session_state['extracted_pdf'] = None
if 'extracted_preview_imgs' not in st.session_state: st.session_state['extracted_preview_imgs'] = []
if 'rotate_states' not in st.session_state: st.session_state['rotate_states'] = {} 
if 'tesseract_path' not in st.session_state: st.session_state['tesseract_path'] = get_local_tesseract_path()
if 'split_results' not in st.session_state: st.session_state['split_results'] = None
if 'global_rot_angle' not in st.session_state: st.session_state['global_rot_angle'] = 0
if 'unlocked_pdf_bytes' not in st.session_state: st.session_state['unlocked_pdf_bytes'] = None
if 'unlocked_file_data' not in st.session_state: st.session_state['unlocked_file_data'] = None

# States for Visual Editors
if 'visual_edit_queue' not in st.session_state: st.session_state['visual_edit_queue'] = []
if 'visual_edit_file_hash' not in st.session_state: st.session_state['visual_edit_file_hash'] = None
if 'visual_sign_queue' not in st.session_state: st.session_state['visual_sign_queue'] = []
if 'visual_sign_file_hash' not in st.session_state: st.session_state['visual_sign_file_hash'] = None

poppler_path = get_local_poppler_path()
tesseract_path = get_local_tesseract_path()

if tesseract_path and HAS_OCR_SUPPORT:
    pytesseract.pytesseract.tesseract_cmd = tesseract_path

# --- SIDEBAR ---
st.sidebar.title("Tools Menu")
category = st.sidebar.selectbox("Choose Category", [
    "Organize & Merge", "Optimize & Repair", "Convert FROM PDF", "Edit & Security"
])

# ==============================================================================
# CATEGORY 1: ORGANIZE & MERGE
# ==============================================================================
if category == "Organize & Merge":
    tool = st.sidebar.radio("Select Tool", ["Merge & Reorder Pages", "Extract Pages", "Split PDF"])

    if tool == "Merge & Reorder Pages":
        st.header("🔗 Merge & Organize Pages")
        st.write("Upload multiple files. Reorder any page from any file.")
        accepted_types = ["pdf", "png", "jpg", "jpeg", "docx", "txt"] if HAS_DOCX_SUPPORT else ["pdf", "png", "jpg", "jpeg", "txt"]
        uploaded_files = st.file_uploader("Add files", type=accepted_types, accept_multiple_files=True)

        if uploaded_files:
            new_files_processed = False
            for file in uploaded_files:
                file_id = f"{file.name}_{file.size}"
                if file_id not in st.session_state['processed_files']:
                    with st.spinner(f"Processing {file.name}..."):
                        temp_pdf_bytes = None
                        if file.type == "application/pdf": temp_pdf_bytes = file.read()
                        elif file.type in ["image/png", "image/jpeg"]:
                            img = Image.open(file)
                            if img.mode == "RGBA": img = img.convert("RGB")
                            b = io.BytesIO(); img.save(b, format='JPEG'); temp_pdf_bytes = img2pdf.convert(b.getvalue())
                        elif file.name.endswith(".txt"): temp_pdf_bytes = text_to_pdf(file.read())
                        elif file.name.endswith(".docx") and HAS_DOCX_SUPPORT:
                            try:
                                t_docx = f"temp_{uuid.uuid4()}.docx"; t_pdf = f"temp_{uuid.uuid4()}.pdf"
                                with open(t_docx, "wb") as f: f.write(file.read())
                                convert_docx(t_docx, t_pdf)
                                with open(t_pdf, "rb") as f: temp_pdf_bytes = f.read()
                                os.remove(t_docx); os.remove(t_pdf)
                            except: st.error(f"Failed to convert {file.name}")

                        if temp_pdf_bytes:
                            reader = PdfReader(io.BytesIO(temp_pdf_bytes))
                            for i, page in enumerate(reader.pages):
                                writer = PdfWriter(); writer.add_page(page); p_bytes = io.BytesIO(); writer.write(p_bytes)
                                st.session_state['page_queue'].append({
                                    'id': str(uuid.uuid4()), 'source': file.name, 'page_num': i + 1,
                                    'bytes': p_bytes.getvalue(), 'rotation': 0
                                })
                            st.session_state['processed_files'].add(file_id)
                            new_files_processed = True
            if new_files_processed: st.rerun()

        if st.session_state['page_queue']:
            total_pages = len(st.session_state['page_queue'])
            st.markdown("---"); c_info, c_clear = st.columns([4, 1]); c_info.info(f"Total Pages: {total_pages}")
            if c_clear.button("Clear All"): st.session_state['page_queue'] = []; st.session_state['processed_files'] = set(); st.rerun()

            with st.expander("🔀 Quick Reorder (Type order)", expanded=False):
                col_ord, col_go = st.columns([4, 1])
                default_order = ", ".join([str(i+1) for i in range(total_pages)])
                with col_ord: new_order = st.text_input("New Order (e.g. 5, 1-4, 6)", value=default_order)
                with col_go:
                    if st.button("Apply"):
                        idxs = parse_order_string(new_order, total_pages)
                        if idxs and len(idxs) > 0: st.session_state['page_queue'] = [st.session_state['page_queue'][i] for i in idxs]; st.rerun()

            st.write("### Page Preview & Reorder")
            cols = st.columns(4)
            for i, item in enumerate(st.session_state['page_queue']):
                with cols[i % 4]:
                    with st.container():
                        st.markdown(f"<div class='page-card'>", unsafe_allow_html=True)
                        st.caption(f"#{i+1} | {item['source']} (Pg {item['page_num']})")
                        thumb = get_page_thumbnail(item['bytes'], poppler_path)
                        rot = item.get('rotation', 0)
                        if thumb and rot != 0: thumb = thumb.rotate(-rot, expand=True)
                        if thumb: st.image(thumb, use_container_width=True)
                        else: st.info("No Preview")
                        
                        c_rot1, c_rot2 = st.columns(2)
                        if c_rot1.button("⟲", key=f"ccw_m_{item['id']}"): st.session_state['page_queue'][i]['rotation'] = (rot - 90) % 360; st.rerun()
                        if c_rot2.button("⟳", key=f"cw_m_{item['id']}"): st.session_state['page_queue'][i]['rotation'] = (rot + 90) % 360; st.rerun()

                        c1, c2, c3 = st.columns([1,1,1])
                        if c1.button("⬅️", key=f"L_{item['id']}") and i > 0: st.session_state['page_queue'][i], st.session_state['page_queue'][i-1] = st.session_state['page_queue'][i-1], st.session_state['page_queue'][i]; st.rerun()
                        if c2.button("❌", key=f"D_{item['id']}"): st.session_state['page_queue'].pop(i); st.rerun()
                        if c3.button("➡️", key=f"R_{item['id']}") and i < total_pages - 1: st.session_state['page_queue'][i], st.session_state['page_queue'][i+1] = st.session_state['page_queue'][i+1], st.session_state['page_queue'][i]; st.rerun()
                        st.markdown("</div>", unsafe_allow_html=True)

            st.markdown("---")
            if st.button("⬇️ Download Final Merged PDF", type="primary"):
                final_merger = PdfWriter()
                for item in st.session_state['page_queue']:
                    r = PdfReader(io.BytesIO(item['bytes'])); page = r.pages[0]
                    if item.get('rotation', 0) != 0: page.rotate(item['rotation'])
                    final_merger.add_page(page)
                out = io.BytesIO(); final_merger.write(out)
                st.download_button("Click to Save PDF", out.getvalue(), "merged_document.pdf", "application/pdf")

    elif tool == "Extract Pages":
        st.header("📄 Extract Pages")
        file = st.file_uploader("Upload PDF", type="pdf")
        if file:
            use_visual = st.checkbox("Enable Visual Page Editor (Rotate/Reorder)", value=False)
            total_pages_source = 0
            
            if use_visual:
                file_hash = f"{file.name}_{file.size}"
                if st.session_state['visual_edit_file_hash'] != file_hash:
                    st.session_state['visual_edit_file_hash'] = file_hash; st.session_state['visual_edit_queue'] = []
                    file.seek(0); reader = PdfReader(file)
                    for i, page in enumerate(reader.pages):
                        writer = PdfWriter(); writer.add_page(page); p_bytes = io.BytesIO(); writer.write(p_bytes)
                        st.session_state['visual_edit_queue'].append({
                            'id': str(uuid.uuid4()), 'page_num': i + 1, 'bytes': p_bytes.getvalue(), 'rotation': 0
                        })
            
                if st.session_state['visual_edit_queue']:
                    total_pages_source = len(st.session_state['visual_edit_queue'])
                    with st.expander("👁️ Visual Editor", expanded=True):
                        cols = st.columns(4)
                        for i, item in enumerate(st.session_state['visual_edit_queue']):
                            with cols[i % 4]:
                                with st.container():
                                    st.markdown(f"<div class='page-card'>", unsafe_allow_html=True)
                                    st.caption(f"Page {i+1}")
                                    thumb = get_page_thumbnail(item['bytes'], poppler_path)
                                    rot = item.get('rotation', 0)
                                    if thumb and rot != 0: thumb = thumb.rotate(-rot, expand=True)
                                    if thumb: st.image(thumb, use_container_width=True)
                                    c_rot1, c_rot2 = st.columns(2)
                                    if c_rot1.button("⟲", key=f"ccw_e_{item['id']}"): st.session_state['visual_edit_queue'][i]['rotation'] = (rot - 90) % 360; st.rerun()
                                    if c_rot2.button("⟳", key=f"cw_e_{item['id']}"): st.session_state['visual_edit_queue'][i]['rotation'] = (rot + 90) % 360; st.rerun()
                                    c1, c2, c3 = st.columns([1,1,1])
                                    if c1.button("⬅️", key=f"L_e_{item['id']}") and i > 0: st.session_state['visual_edit_queue'][i], st.session_state['visual_edit_queue'][i-1] = st.session_state['visual_edit_queue'][i-1], st.session_state['visual_edit_queue'][i]; st.rerun()
                                    if c3.button("➡️", key=f"R_e_{item['id']}") and i < total_pages_source - 1: st.session_state['visual_edit_queue'][i], st.session_state['visual_edit_queue'][i+1] = st.session_state['visual_edit_queue'][i+1], st.session_state['visual_edit_queue'][i]; st.rerun()
                                    st.markdown("</div>", unsafe_allow_html=True)
            else:
                file.seek(0); reader_check = PdfReader(file); total_pages_source = len(reader_check.pages)
                
            st.markdown("---")
            page_input = st.text_input("Pages to Extract (e.g. 1, 3-5)", "1")
            if st.button("Preview & Process"):
                try:
                    idxs = parse_order_string(page_input, total_pages_source)
                    if idxs:
                        writer = PdfWriter(); preview_imgs = []
                        if use_visual:
                             for i, page_idx in enumerate(idxs):
                                item = st.session_state['visual_edit_queue'][page_idx]
                                p_reader = PdfReader(io.BytesIO(item['bytes'])); page = p_reader.pages[0]
                                if item.get('rotation', 0) != 0: page.rotate(item['rotation'])
                                writer.add_page(page)
                                if poppler_path:
                                    temp_writer = PdfWriter(); temp_writer.add_page(page); temp_bytes = io.BytesIO(); temp_writer.write(temp_bytes)
                                    thumb = get_page_thumbnail(temp_bytes.getvalue(), poppler_path)
                                    if thumb: preview_imgs.append((i+1, thumb))
                        else:
                            file.seek(0); reader = PdfReader(file)
                            for i, page_idx in enumerate(idxs):
                                page = reader.pages[page_idx]; writer.add_page(page)
                                if poppler_path:
                                     temp_writer = PdfWriter(); temp_writer.add_page(page); temp_bytes = io.BytesIO(); temp_writer.write(temp_bytes)
                                     thumb = get_page_thumbnail(temp_bytes.getvalue(), poppler_path)
                                     if thumb: preview_imgs.append((i+1, thumb))
                        out = io.BytesIO(); writer.write(out)
                        st.session_state['extracted_pdf'] = out.getvalue()
                        st.session_state['extracted_preview_imgs'] = preview_imgs
                    else: st.error("No valid pages selected.")
                except Exception as e: st.error(f"Error processing pages: {e}")

            if st.session_state.get('extracted_pdf'):
                st.markdown("### Result Preview")
                if st.session_state['extracted_preview_imgs']:
                    cols = st.columns(4)
                    for i, (num, img) in enumerate(st.session_state['extracted_preview_imgs']):
                        with cols[i % 4]: st.image(img, caption=f"New Page {num}", use_container_width=True)
                st.download_button(label="⬇️ Download Extracted PDF", data=st.session_state['extracted_pdf'], file_name="extracted.pdf", mime="application/pdf", type="primary")

    elif tool == "Split PDF":
        st.header("✂️ Split PDF")
        mode = st.radio("Split Mode", ["Custom Ranges", "Fixed Page Range", "Split into N Files", "Extract All Pages"])
        file = st.file_uploader("Upload PDF", type="pdf")
        if file:
            use_visual = st.checkbox("Enable Visual Page Editor (Rotate/Reorder)", value=False)
            total_pages_source = 0
            if use_visual:
                file_hash = f"{file.name}_{file.size}"
                if st.session_state['visual_edit_file_hash'] != file_hash:
                    st.session_state['visual_edit_file_hash'] = file_hash; st.session_state['visual_edit_queue'] = []
                    file.seek(0); reader = PdfReader(file)
                    for i, page in enumerate(reader.pages):
                        writer = PdfWriter(); writer.add_page(page); p_bytes = io.BytesIO(); writer.write(p_bytes)
                        st.session_state['visual_edit_queue'].append({
                            'id': str(uuid.uuid4()), 'page_num': i + 1, 'bytes': p_bytes.getvalue(), 'rotation': 0
                        })
                
                if st.session_state['visual_edit_queue']:
                    total_pages_source = len(st.session_state['visual_edit_queue'])
                    with st.expander("👁️ Visual Editor", expanded=True):
                        cols = st.columns(4)
                        for i, item in enumerate(st.session_state['visual_edit_queue']):
                            with cols[i % 4]:
                                with st.container():
                                    st.markdown(f"<div class='page-card'>", unsafe_allow_html=True)
                                    st.caption(f"Page {i+1}")
                                    thumb = get_page_thumbnail(item['bytes'], poppler_path)
                                    rot = item.get('rotation', 0)
                                    if thumb and rot != 0: thumb = thumb.rotate(-rot, expand=True)
                                    if thumb: st.image(thumb, use_container_width=True)
                                    c_rot1, c_rot2 = st.columns(2)
                                    if c_rot1.button("⟲", key=f"ccw_s_{item['id']}"): st.session_state['visual_edit_queue'][i]['rotation'] = (rot - 90) % 360; st.rerun()
                                    if c_rot2.button("⟳", key=f"cw_s_{item['id']}"): st.session_state['visual_edit_queue'][i]['rotation'] = (rot + 90) % 360; st.rerun()
                                    c1, c2, c3 = st.columns([1,1,1])
                                    if c1.button("⬅️", key=f"L_s_{item['id']}") and i > 0: st.session_state['visual_edit_queue'][i], st.session_state['visual_edit_queue'][i-1] = st.session_state['visual_edit_queue'][i-1], st.session_state['visual_edit_queue'][i]; st.rerun()
                                    if c3.button("➡️", key=f"R_s_{item['id']}") and i < total_pages_source - 1: st.session_state['visual_edit_queue'][i], st.session_state['visual_edit_queue'][i+1] = st.session_state['visual_edit_queue'][i+1], st.session_state['visual_edit_queue'][i]; st.rerun()
                                    st.markdown("</div>", unsafe_allow_html=True)
            else:
                file.seek(0); reader = PdfReader(file); total_pages_source = len(reader.pages)
            
            st.markdown("---"); st.info(f"Total Pages: {total_pages_source}")
            split_groups = []
            
            if mode == "Custom Ranges":
                range_str = st.text_input("Ranges (comma separated)", "1-5, 6-10") 
                if range_str:
                    try:
                        parts = [p.strip() for p in range_str.split(',') if p.strip()]
                        for p in parts:
                            if '-' in p: s, e = map(int, p.split('-')); split_groups.append(list(range(s-1, e)))
                            else: split_groups.append([int(p)-1])
                    except: st.error("Invalid range format.")
            elif mode.startswith("Fixed Page Range"):
                chunk_size = st.number_input("Pages per file:", min_value=1, max_value=total_pages_source, value=1)
                for i in range(0, total_pages_source, chunk_size): split_groups.append(list(range(i, min(i + chunk_size, total_pages_source))))
            elif mode.startswith("Split into N Files"):
                num_files = st.number_input("Number of files:", min_value=2, max_value=total_pages_source, value=2)
                chunk_size = math.ceil(total_pages_source / num_files)
                for i in range(0, total_pages_source, chunk_size): split_groups.append(list(range(i, min(i + chunk_size, total_pages_source))))
            elif mode == "Extract All Pages":
                split_groups = [[i] for i in range(total_pages_source)]

            if st.button("Process Split", type="primary"):
                if not split_groups: st.error("No ranges defined.")
                else:
                    try:
                        files_data = []; zip_buffer = io.BytesIO()
                        with zipfile.ZipFile(zip_buffer, "w") as zf:
                            for idx, pages in enumerate(split_groups):
                                writer = PdfWriter(); valid = False
                                for p_idx in pages:
                                    if 0 <= p_idx < total_pages_source:
                                        page = None
                                        if use_visual:
                                             item = st.session_state['visual_edit_queue'][p_idx]
                                             p_reader = PdfReader(io.BytesIO(item['bytes'])); page = p_reader.pages[0]
                                             if item.get('rotation', 0) != 0: page.rotate(item['rotation'])
                                        else:
                                             file.seek(0); safe_reader = PdfReader(file); page = safe_reader.pages[p_idx]
                                        writer.add_page(page); valid = True
                                if valid:
                                    pdf_bytes = io.BytesIO(); writer.write(pdf_bytes); data = pdf_bytes.getvalue()
                                    if len(pages) == 1: name = f"Page_{pages[0]+1}.pdf"
                                    else: name = f"Split_{idx+1}_(Pg{pages[0]+1}-{pages[-1]+1}).pdf"
                                    zf.writestr(name, data); files_data.append((name, data))
                        st.session_state['split_results'] = {'zip': zip_buffer.getvalue(), 'files': files_data}
                    except Exception as e: st.error(f"Error splitting PDF: {e}")

            if st.session_state['split_results']:
                st.success("Split Complete!")
                res = st.session_state['split_results']
                st.download_button("⬇️ Download All (ZIP)", res['zip'], "split_files.zip", "application/zip", type="primary")
                st.markdown("### Individual Downloads")
                for name, data in res['files']: st.download_button(f"📄 {name}", data, name, "application/pdf")

# ==============================================================================
# CATEGORY 2: OPTIMIZE & REPAIR
# ==============================================================================
elif category == "Optimize & Repair":
    tool = st.sidebar.radio("Select Tool", ["Compress PDF", "Repair PDF"])

    # --- 1. COMPRESS PDF ---
    if tool == "Compress PDF":
        st.header("📉 Compress PDF")
        file = st.file_uploader("Upload PDF", type="pdf")
        if file: file.seek(0, os.SEEK_END); orig_size = file.tell(); file.seek(0); st.info(f"Original File Size: {orig_size/1024:.2f} KB")
        comp_mode = st.radio("Compression Level", ["Basic (Lossless)", "Strong (Flatten to Images)"])
        quality_val = 70
        if "Strong" in comp_mode: quality_val = st.slider("Compression Strength (Image Quality)", min_value=10, max_value=95, value=60)
        if file and st.button("Compress"):
            file.seek(0)
            if comp_mode.startswith("Basic"):
                try:
                    with open("temp_in.pdf", "wb") as f: f.write(file.read())
                    with pikepdf.open("temp_in.pdf") as pdf: pdf.save("temp_out.pdf", compress_streams=True, object_stream_mode=pikepdf.ObjectStreamMode.generate)
                    new_size = os.path.getsize("temp_out.pdf")
                    st.success(f"Compressed! New Size: {new_size/1024:.2f} KB")
                    with open("temp_out.pdf", "rb") as f: st.download_button("Download Compressed PDF", f.read(), "compressed_lossless.pdf", "application/pdf")
                    os.remove("temp_in.pdf"); os.remove("temp_out.pdf")
                except Exception as e: st.error(f"Error: {e}")
            else:
                st.info(f"Converting pages to images ({quality_val}% Quality JPEG) and rebuilding PDF...")
                try:
                    common_args = {"dpi": 150}
                    if poppler_path: images = convert_from_bytes(file.read(), poppler_path=poppler_path, **common_args)
                    else: images = convert_from_bytes(file.read(), **common_args)
                    img_bytes_list = []
                    for img in images: b = io.BytesIO(); img.save(b, format='JPEG', quality=quality_val); img_bytes_list.append(b.getvalue())
                    pdf_bytes = img2pdf.convert(img_bytes_list)
                    new_size = len(pdf_bytes)
                    st.success(f"Done! New Size: {new_size/1024:.2f} KB")
                    st.download_button("Download Compressed PDF", pdf_bytes, "compressed_strong.pdf", "application/pdf")
                except Exception as e: st.error(f"Error: {e}")

    # --- 2. REPAIR PDF ---
    elif tool == "Repair PDF":
        st.header("🔧 Repair PDF")
        file = st.file_uploader("Upload Corrupted PDF", type="pdf")
        if file and st.button("Repair & Download"):
            try:
                with open("corrupt.pdf", "wb") as f: f.write(file.read())
                with pikepdf.open("corrupt.pdf", allow_overwriting_input=True) as pdf: pdf.save("repaired.pdf")
                with open("repaired.pdf", "rb") as f:
                    st.success("File processed.")
                    st.download_button("Download Repaired PDF", f.read(), "repaired.pdf", "application/pdf")
                os.remove("corrupt.pdf"); os.remove("repaired.pdf")
            except Exception as e: st.error(f"Repair failed: {e}. The file might be too damaged.")

# ==============================================================================
# CATEGORY 3: CONVERT FROM PDF
# ==============================================================================
elif category == "Convert FROM PDF":
    tool = st.sidebar.radio("Select Tool", ["PDF to Images", "PDF to Word", "PDF to Excel", "PDF to Text", "PDF to PowerPoint", "OCR PDF (Searchable)"])

    if tool == "PDF to Images":
        st.header("🖼️ PDF to Images")
        file = st.file_uploader("Upload PDF", type="pdf")
        col_set, _ = st.columns([2,1])
        with col_set: quality_setting = st.select_slider("Conversion Speed vs Quality", options=["Screen (72 dpi)", "Standard (150 dpi)", "Print (300 dpi)"], value="Standard (150 dpi)")
        dpi_map = {"Screen (72 dpi)": 72, "Standard (150 dpi)": 150, "Print (300 dpi)": 300}; selected_dpi = dpi_map[quality_setting]
        if file and st.button("Convert to Images"):
            common_args = {"dpi": selected_dpi, "thread_count": 4}
            if poppler_path: images = convert_from_bytes(file.read(), poppler_path=poppler_path, **common_args)
            else: images = convert_from_bytes(file.read(), **common_args)
            zip_buffer = io.BytesIO()
            with zipfile.ZipFile(zip_buffer, "w") as zf:
                for i, img in enumerate(images):
                    c1, c2 = st.columns([1,3])
                    with c1: st.image(img, use_container_width=True)
                    b = io.BytesIO(); img.save(b, 'JPEG'); img_bytes = b.getvalue()
                    zf.writestr(f"page_{i+1}.jpg", img_bytes)
                    with c2: st.download_button(f"Download Page {i+1}", img_bytes, f"p{i+1}.jpg", "image/jpeg")
            st.markdown("---")
            st.download_button(label="⬇️ Download All (ZIP)", data=zip_buffer.getvalue(), file_name="all_images.zip", mime="application/zip", type="primary")

    elif tool == "PDF to Word":
        st.header("📝 PDF to Word (.docx)")
        file = st.file_uploader("Upload PDF", type="pdf")
        if file and st.button("Convert to Word"):
            try:
                temp_pdf = f"temp_{uuid.uuid4()}.pdf"; temp_docx = f"temp_{uuid.uuid4()}.docx"
                with open(temp_pdf, "wb") as f: f.write(file.read())
                cv = Converter(temp_pdf); cv.convert(temp_docx, start=0, end=None); cv.close()
                with open(temp_docx, "rb") as f: docx_bytes = f.read()
                st.success("Success!"); st.download_button("Download Word Doc", docx_bytes, "converted.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document", type="primary")
                os.remove(temp_pdf); os.remove(temp_docx)
            except Exception as e: st.error(f"Error: {e}")

    elif tool == "PDF to Excel":
        st.header("📊 PDF to Excel (.xlsx)")
        file = st.file_uploader("Upload PDF", type="pdf")
        if file and st.button("Convert to Excel"):
            try:
                all_tables = []
                with pdfplumber.open(file) as pdf:
                    for i, page in enumerate(pdf.pages):
                        tables = page.extract_tables()
                        for table in tables: df = pd.DataFrame(table); all_tables.append(df)
                if all_tables:
                    output = io.BytesIO()
                    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
                        for i, df in enumerate(all_tables): df.to_excel(writer, sheet_name=f"Table_{i+1}", index=False, header=False)
                    st.success(f"Found {len(all_tables)} tables!"); st.download_button("Download Excel File", output.getvalue(), "extracted_tables.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", type="primary")
                else: st.warning("No tables found.")
            except Exception as e: st.error(f"Error: {e}")

    elif tool == "PDF to Text":
        st.header("📄 PDF to Text (.txt)")
        file = st.file_uploader("Upload PDF", type="pdf")
        if file and st.button("Extract Text"):
            try:
                reader = PdfReader(file); full_text = ""
                for page in reader.pages: full_text += page.extract_text() + "\n\n"
                st.success("Done!"); st.text_area("Preview", full_text[:500] + "...", height=200)
                st.download_button("Download Text File", full_text, "extracted_text.txt", "text/plain", type="primary")
            except Exception as e: st.error(f"Error: {e}")

    elif tool == "PDF to PowerPoint":
        st.header("📽️ PDF to PowerPoint (.pptx)")
        if not HAS_PPTX_SUPPORT: st.error("⚠️ `python-pptx` library missing.")
        else:
            mode = st.radio("Conversion Mode", ["Image-based (Better Layout)", "Editable (Text Boxes)"])
            file = st.file_uploader("Upload PDF", type="pdf")
            if file and st.button("Convert to PPTX"):
                try:
                    if mode.startswith("Editable"):
                        with st.spinner("Analyzing text layout (Editable Mode)..."):
                            pptx_bytes = create_editable_pptx(file)
                            st.success("Editable Conversion Complete!")
                            st.download_button("⬇️ Download PPTX", pptx_bytes, "editable_presentation.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation", type="primary")
                    else:
                        with st.spinner("Converting pages to slides (Image Mode)..."):
                            images = convert_from_bytes(file.read(), poppler_path=poppler_path, dpi=150)
                            prs = Presentation(); blank_slide_layout = prs.slide_layouts[6]
                            if images:
                                width_px, height_px = images[0].size; aspect_ratio = width_px / height_px
                                prs.slide_width = Inches(10); prs.slide_height = Inches(10 / aspect_ratio)
                            for img in images:
                                slide = prs.slides.add_slide(blank_slide_layout); img_stream = io.BytesIO(); img.save(img_stream, format="PNG"); slide.shapes.add_picture(img_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
                            pptx_out = io.BytesIO(); prs.save(pptx_out)
                            st.success("Image Conversion Complete!")
                            st.download_button("⬇️ Download PPTX", pptx_out.getvalue(), "presentation_images.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation", type="primary")
                except Exception as e: st.error(f"Error converting to PowerPoint: {e}")

    elif tool == "OCR PDF (Searchable)":
        st.header("🔍 OCR PDF (Searchable)")
        if not HAS_OCR_SUPPORT: st.error("❌ Python library `pytesseract` is missing.")
        else:
            if not st.session_state.get('tesseract_path'):
                st.warning("⚠️ Tesseract.exe not found automatically.")
                manual_tesseract = st.text_input("Paste path manually:", value=r"C:\Users\Surface\OneDrive\Desktop\Add-in project\tesseract.exe")
                if os.path.exists(manual_tesseract): st.session_state['tesseract_path'] = manual_tesseract; pytesseract.pytesseract.tesseract_cmd = manual_tesseract; st.success("✅ Tesseract found!"); st.rerun()
            if st.session_state.get('tesseract_path') and os.path.exists(st.session_state['tesseract_path']):
                file = st.file_uploader("Upload Scanned PDF or Image", type=["pdf", "png", "jpg", "jpeg"])
                lang = st.selectbox("Language", ["eng", "spa", "fra", "deu"])
                if file and st.button("Run OCR", type="primary"):
                    try:
                        def process_ocr_page(image):
                            return pytesseract.image_to_pdf_or_hocr(image, extension='pdf', lang=lang)

                        with st.spinner("Running OCR in parallel..."):
                            if file.name.endswith(".pdf"): images = convert_from_bytes(file.read(), poppler_path=poppler_path, dpi=200)
                            else: images = [Image.open(file)]
                            
                            pdf_writer = PdfWriter()
                            # Use ThreadPoolExecutor to process pages concurrently
                            with concurrent.futures.ThreadPoolExecutor(max_workers=4) as executor:
                                results = list(executor.map(process_ocr_page, images))
                            
                            for page_bytes in results:
                                pdf_reader = PdfReader(io.BytesIO(page_bytes))
                                pdf_writer.add_page(pdf_reader.pages[0])
                                
                            out = io.BytesIO(); pdf_writer.write(out)
                            st.success(f"OCR Complete! Processed {len(images)} pages."); st.download_button("Download Searchable PDF", out.getvalue(), "ocr_searchable.pdf", "application/pdf")
                    except Exception as e: st.error(f"OCR Error: {e}")

# ==============================================================================
# CATEGORY 4: EDIT & SECURITY
# ==============================================================================
elif category == "Edit & Security":
    tool = st.sidebar.radio("Select Tool", ["Add Watermark", "Add Page Numbers", "Header & Footer", "Rotate PDF", "Crop PDF", "Sign PDF", "Lock PDF", "Decrypt / Unlock PDF"])

    if tool == "Add Watermark":
        st.header("💧 Add Watermark")
        file = st.file_uploader("Upload PDF", type="pdf")
        c1, c2, c3 = st.columns(3)
        with c1: wm_text = st.text_input("Watermark Text", "CONFIDENTIAL")
        with c2: wm_style = st.selectbox("Style", ["Single (Centered)", "Tiled (Repeat)"])
        with c3:
            if "Single" in wm_style: wm_pos = st.selectbox("Position", ["Center", "Top Left", "Top Center", "Top Right", "Bottom Left", "Bottom Center", "Bottom Right", "Custom (Manual X/Y)"])
        if "Tiled" in wm_style:
            c_sp1, c_sp2 = st.columns(2)
            with c_sp1: gap_x = st.slider("Horizontal Spacing", 50, 500, 200)
            with c_sp2: gap_y = st.slider("Vertical Spacing", 50, 500, 200)
        c4, c5, c6 = st.columns(3)
        with c4: font_size = st.slider("Font Size", 10, 100, 50)
        with c5: opacity = st.slider("Opacity", 0.1, 1.0, 0.5)
        with c6: rotation = st.slider("Rotation", 0, 360, 45)
        c7, c8 = st.columns([1, 1])
        with c7: font_color = st.color_picker("Color", "#808080")
        with c8: font_face = font_selector_component("wm")
        
        # Custom X/Y Inputs
        custom_x, custom_y = 100, 100
        # Determine Page Size for Sliders
        max_w, max_h = 612, 792 # Defaults (Letter)
        if file:
            try:
                # Quick peek for dimensions
                r_temp = PdfReader(file)
                p_temp = r_temp.pages[0]
                max_w = int(float(p_temp.mediabox.width))
                max_h = int(float(p_temp.mediabox.height))
                file.seek(0)
            except: pass

        if "Single" in wm_style and wm_pos == "Custom (Manual X/Y)":
            c_cust1, c_cust2 = st.columns(2)
            with c_cust1: custom_x = st.slider("X Coordinate", 0, max_w, int(max_w/2))
            with c_cust2: custom_y = st.slider("Y Coordinate", 0, max_h, int(max_h/2))

        if file:
            try:
                st.markdown("### Live Preview")
                reader = PdfReader(file); page_1 = reader.pages[0]; pg_w = float(page_1.mediabox.width); pg_h = float(page_1.mediabox.height)
                packet = io.BytesIO(); c = canvas.Canvas(packet, pagesize=(pg_w, pg_h))
                hex_color = font_color.lstrip('#'); rgb = tuple(int(hex_color[i:i+2], 16)/255.0 for i in (0, 2, 4))
                c.setFont(font_face, font_size); c.setFillColorRGB(*rgb, alpha=opacity); c.saveState()
                if "Tiled" in wm_style:
                    c.rotate(rotation); actual_gap_x = gap_x + (len(wm_text) * 2); actual_gap_y = gap_y
                    for x in range(-int(pg_w*2), int(pg_w*2), actual_gap_x):
                        for y in range(-int(pg_h*2), int(pg_h*2), actual_gap_y): c.drawString(x, y, wm_text)
                else:
                    x, y = 0, 0; margin = 50
                    if wm_pos == "Custom (Manual X/Y)": x, y = custom_x, custom_y
                    elif "Center" in wm_pos: x = pg_w/2
                    elif "Left" in wm_pos: x = margin
                    elif "Right" in wm_pos: x = pg_w - margin
                    
                    if wm_pos != "Custom (Manual X/Y)":
                        if "Center" in wm_pos and "Top" not in wm_pos and "Bottom" not in wm_pos: y = pg_h/2
                        elif "Top" in wm_pos: y = pg_h - margin
                        elif "Bottom" in wm_pos: y = margin
                    
                    c.translate(x, y); c.rotate(rotation); c.drawCentredString(0, 0, wm_text)
                c.restoreState(); c.save(); packet.seek(0)
                wm_page = PdfReader(packet).pages[0]; first_page = page_1; first_page.merge_page(wm_page)
                writer = PdfWriter(); writer.add_page(first_page); temp_out = io.BytesIO(); writer.write(temp_out)
                thumb = get_page_thumbnail(temp_out.getvalue(), poppler_path, width=800)
                if thumb: st.image(thumb, width=500)
                file.seek(0)
            except Exception as e: st.error(f"Preview Error: {e}")
            if st.button("Apply to All Pages & Download", type="primary"):
                try:
                    reader = PdfReader(file); writer = PdfWriter(); hex_color = font_color.lstrip('#'); rgb = tuple(int(hex_color[i:i+2], 16)/255.0 for i in (0, 2, 4))
                    for page in reader.pages:
                        pg_w = float(page.mediabox.width); pg_h = float(page.mediabox.height); packet = io.BytesIO(); c = canvas.Canvas(packet, pagesize=(pg_w, pg_h))
                        c.setFont(font_face, font_size); c.setFillColorRGB(*rgb, alpha=opacity); c.saveState()
                        if "Tiled" in wm_style:
                            c.rotate(rotation); actual_gap_x = gap_x + (len(wm_text) * 2); actual_gap_y = gap_y
                            for x in range(-int(pg_w*2), int(pg_w*2), actual_gap_x):
                                for y in range(-int(pg_h*2), int(pg_h*2), actual_gap_y): c.drawString(x, y, wm_text)
                        else:
                            x, y = 0, 0; margin = 50
                            if wm_pos == "Custom (Manual X/Y)": x, y = custom_x, custom_y
                            elif "Center" in wm_pos: x = pg_w/2
                            elif "Left" in wm_pos: x = margin
                            elif "Right" in wm_pos: x = pg_w - margin
                            if wm_pos != "Custom (Manual X/Y)":
                                if "Center" in wm_pos and "Top" not in wm_pos and "Bottom" not in wm_pos: y = pg_h/2
                                elif "Top" in wm_pos: y = pg_h - margin
                                elif "Bottom" in wm_pos: y = margin
                            c.translate(x, y); c.rotate(rotation); c.drawCentredString(0, 0, wm_text)
                        c.restoreState(); c.save(); packet.seek(0)
                        wm_layer = PdfReader(packet).pages[0]; page.merge_page(wm_layer); writer.add_page(page)
                    out = io.BytesIO(); writer.write(out); st.download_button("Download PDF", out.getvalue(), "watermarked.pdf", "application/pdf")
                except Exception as e: st.error(f"Error: {e}")

    elif tool == "Add Page Numbers":
        st.header("🔢 Add Page Numbers")
        file = st.file_uploader("Upload PDF", type="pdf")
        c1, c2, c3 = st.columns(3)
        with c1: position = st.selectbox("Position", ["Bottom Center", "Bottom Right", "Bottom Left", "Top Center", "Top Right", "Top Left"])
        with c2: style_fmt = st.selectbox("Format", ["Page 1", "1", "Page 1 of N", "1 of N"])
        with c3: opacity = st.slider("Opacity", 0.1, 1.0, 1.0)
        c4, c5, c6 = st.columns(3)
        with c4: font_size = st.number_input("Size", 8, 72, 12)
        with c5: font_color = st.color_picker("Color", "#000000")
        with c6: font_face = font_selector_component("pnum")

        if file:
            try:
                st.markdown("### Live Preview (Page 1)")
                packet = io.BytesIO(); reader = PdfReader(file); first_page = reader.pages[0]; pg_width = float(first_page.mediabox.width); pg_height = float(first_page.mediabox.height)
                c = canvas.Canvas(packet, pagesize=(pg_width, pg_height)); hex_color = font_color.lstrip('#'); rgb = tuple(int(hex_color[i:i+2], 16)/255.0 for i in (0, 2, 4))
                c.setFillColorRGB(*rgb, alpha=opacity); c.setFont(font_face, font_size)
                text_str = style_fmt.replace("1", "1").replace("N", "5"); margin = 20
                if "Bottom" in position: y = margin
                else: y = pg_height - margin
                if "Left" in position: x = margin; draw_func = c.drawString
                elif "Right" in position: x = pg_width - margin; draw_func = c.drawRightString
                else: x = pg_width / 2; draw_func = c.drawCentredString
                draw_func(x, y, text_str); c.save(); packet.seek(0)
                num_page = PdfReader(packet).pages[0]; first_page.merge_page(num_page)
                writer = PdfWriter(); writer.add_page(first_page); temp = io.BytesIO(); writer.write(temp)
                thumb = get_page_thumbnail(temp.getvalue(), poppler_path, width=800)
                if thumb: st.image(thumb, width=400)
                file.seek(0)
            except Exception as e: st.error(f"Preview Error: {e}")
            if st.button("Apply & Download", type="primary"):
                try:
                    reader = PdfReader(file); writer = PdfWriter(); total_pages = len(reader.pages)
                    for i, page in enumerate(reader.pages):
                        packet = io.BytesIO(); pg_width = float(page.mediabox.width); pg_height = float(page.mediabox.height); c = canvas.Canvas(packet, pagesize=(pg_width, pg_height))
                        hex_color = font_color.lstrip('#'); rgb = tuple(int(hex_color[i:i+2], 16)/255.0 for i in (0, 2, 4))
                        c.setFillColorRGB(*rgb, alpha=opacity); c.setFont(font_face, font_size)
                        p_num = i + 1; text_str = style_fmt.replace("1", str(p_num)).replace("N", str(total_pages)); margin = 20
                        if "Bottom" in position: y = margin
                        else: y = pg_height - margin
                        if "Left" in position: x = margin; draw_func = c.drawString
                        elif "Right" in position: x = pg_width - margin; draw_func = c.drawRightString
                        else: x = pg_width / 2; draw_func = c.drawCentredString
                        draw_func(x, y, text_str); c.save(); packet.seek(0)
                        num_layer = PdfReader(packet).pages[0]; page.merge_page(num_layer); writer.add_page(page)
                    out = io.BytesIO(); writer.write(out); st.download_button("Download PDF", out.getvalue(), "numbered.pdf", "application/pdf")
                except Exception as e: st.error(f"Error: {e}")

    elif tool == "Header & Footer":
        st.header("📑 Header & Footer")
        file = st.file_uploader("Upload PDF", type="pdf")
        c1, c2, c3 = st.columns(3)
        with c1: position = st.selectbox("Position", ["Top Left", "Top Center", "Top Right", "Bottom Left", "Bottom Center", "Bottom Right"])
        with c2: user_text = st.text_input("Text Content", "My Document")
        with c3: opacity = st.slider("Opacity", 0.1, 1.0, 1.0)
        c4, c5, c6 = st.columns(3)
        with c4: font_size = st.number_input("Size", 8, 72, 12)
        with c5: font_color = st.color_picker("Color", "#000000")
        with c6: font_face = font_selector_component("hf")

        if file:
            try:
                st.markdown("### Live Preview (Page 1)")
                packet = io.BytesIO(); reader = PdfReader(file); first_page = reader.pages[0]; pg_width = float(first_page.mediabox.width); pg_height = float(first_page.mediabox.height)
                c = canvas.Canvas(packet, pagesize=(pg_width, pg_height)); hex_color = font_color.lstrip('#'); rgb = tuple(int(hex_color[i:i+2], 16)/255.0 for i in (0, 2, 4))
                c.setFillColorRGB(*rgb, alpha=opacity); c.setFont(font_face, font_size)
                margin = 20
                if "Bottom" in position: y = margin
                else: y = pg_height - margin
                if "Left" in position: x = margin; draw_func = c.drawString
                elif "Right" in position: x = pg_width - margin; draw_func = c.drawRightString
                else: x = pg_width / 2; draw_func = c.drawCentredString
                draw_func(x, y, user_text); c.save(); packet.seek(0)
                hf_page = PdfReader(packet).pages[0]; first_page.merge_page(hf_page)
                writer = PdfWriter(); writer.add_page(first_page); temp = io.BytesIO(); writer.write(temp)
                thumb = get_page_thumbnail(temp.getvalue(), poppler_path, width=800)
                if thumb: st.image(thumb, width=400)
                file.seek(0)
            except Exception as e: st.error(f"Preview Error: {e}")
            if st.button("Apply & Download", type="primary"):
                try:
                    reader = PdfReader(file); writer = PdfWriter()
                    for i, page in enumerate(reader.pages):
                        packet = io.BytesIO(); pg_width = float(page.mediabox.width); pg_height = float(page.mediabox.height); c = canvas.Canvas(packet, pagesize=(pg_width, pg_height))
                        hex_color = font_color.lstrip('#'); rgb = tuple(int(hex_color[i:i+2], 16)/255.0 for i in (0, 2, 4))
                        c.setFillColorRGB(*rgb, alpha=opacity); c.setFont(font_face, font_size)
                        margin = 20
                        if "Bottom" in position: y = margin
                        else: y = pg_height - margin
                        if "Left" in position: x = margin; draw_func = c.drawString
                        elif "Right" in position: x = pg_width - margin; draw_func = c.drawRightString
                        else: x = pg_width / 2; draw_func = c.drawCentredString
                        draw_func(x, y, user_text); c.save(); packet.seek(0)
                        hf_layer = PdfReader(packet).pages[0]; page.merge_page(hf_layer); writer.add_page(page)
                    out = io.BytesIO(); writer.write(out); st.download_button("Download PDF", out.getvalue(), "document_with_header.pdf", "application/pdf")
                except Exception as e: st.error(f"Error: {e}")

    elif tool == "Rotate PDF":
        st.header("🔄 Rotate PDF")
        mode = st.radio("Rotate Mode", ["Rotate All Pages", "Rotate Individual Pages"])
        file = st.file_uploader("Upload PDF", type="pdf")
        if file:
            if mode == "Rotate All Pages":
                if 'global_rot_angle' not in st.session_state: st.session_state['global_rot_angle'] = 0
                col1, col2, col3 = st.columns([1,1,2])
                with col1:
                    if st.button("⟲ Left"): st.session_state['global_rot_angle'] = (st.session_state['global_rot_angle'] - 90) % 360
                with col2:
                    if st.button("⟳ Right"): st.session_state['global_rot_angle'] = (st.session_state['global_rot_angle'] + 90) % 360
                with col3: st.write(f"**Current Rotation:** {st.session_state['global_rot_angle']}°")
                rot = st.session_state['global_rot_angle']
                try:
                    reader = PdfReader(file); p1 = reader.pages[0]; p1.rotate(rot)
                    writer = PdfWriter(); writer.add_page(p1); temp = io.BytesIO(); writer.write(temp)
                    st.write("### Preview (Page 1)")
                    thumb = get_page_thumbnail(temp.getvalue(), poppler_path, width=400)
                    if thumb: st.image(thumb, width=300)
                    file.seek(0)
                except: st.error("Preview failed")
                if st.button("Rotate All & Download", type="primary"):
                    reader = PdfReader(file); writer = PdfWriter()
                    for page in reader.pages: page.rotate(rot); writer.add_page(page)
                    out = io.BytesIO(); writer.write(out)
                    st.download_button("Download Rotated PDF", out.getvalue(), "rotated_all.pdf", "application/pdf")
            else:
                file_id = f"{file.name}_{file.size}_rot"
                if 'current_rot_file' not in st.session_state or st.session_state['current_rot_file'] != file_id:
                    st.session_state['current_rot_file'] = file_id; st.session_state['rotate_states'] = {} 
                reader = PdfReader(file)
                total_pages = len(reader.pages)
                st.write(f"Total Pages: {total_pages}")
                cols = st.columns(4)
                for i in range(total_pages):
                    current_angle = st.session_state['rotate_states'].get(i, 0)
                    with cols[i % 4]:
                        with st.container():
                            st.markdown(f"<div class='page-card'>", unsafe_allow_html=True)
                            st.caption(f"Page {i+1}")
                            writer = PdfWriter(); writer.add_page(reader.pages[i]); p_bytes = io.BytesIO(); writer.write(p_bytes)
                            thumb = get_page_thumbnail(p_bytes.getvalue(), poppler_path)
                            if thumb:
                                rotated_thumb = thumb.rotate(-current_angle, expand=True)
                                st.image(rotated_thumb, use_container_width=True)
                            else: st.info("No Preview")
                            c1, c2 = st.columns(2)
                            if c1.button("⟲", key=f"ccw_{i}"): st.session_state['rotate_states'][i] = (current_angle - 90) % 360; st.rerun()
                            if c2.button("⟳", key=f"cw_{i}"): st.session_state['rotate_states'][i] = (current_angle + 90) % 360; st.rerun()
                            st.caption(f"Rot: {st.session_state['rotate_states'].get(i, 0)}°")
                            st.markdown("</div>", unsafe_allow_html=True)
                st.markdown("---")
                if st.button("Apply Rotations & Download", type="primary"):
                    final_writer = PdfWriter(); file.seek(0); r = PdfReader(file)
                    for i, page in enumerate(r.pages):
                        angle = st.session_state['rotate_states'].get(i, 0)
                        if angle != 0: page.rotate(angle)
                        final_writer.add_page(page)
                    out = io.BytesIO(); final_writer.write(out)
                    st.download_button("Download Result", out.getvalue(), "individual_rotated.pdf", "application/pdf")

    elif tool == "Crop PDF":
        st.header("✂️ Crop PDF")
        file = st.file_uploader("Upload PDF", type="pdf")
        st.write("Adjust crop margins (in points). 72 points = 1 inch.")
        c1, c2 = st.columns(2)
        with c1: left = st.slider("Left Margin", 0, 200, 0)
        with c2: right = st.slider("Right Margin", 0, 200, 0)
        c3, c4 = st.columns(2)
        with c3: top = st.slider("Top Margin", 0, 200, 0)
        with c4: bottom = st.slider("Bottom Margin", 0, 200, 0)
        if file:
            try:
                reader = PdfReader(file); page = reader.pages[0]; orig_ur = page.mediabox.upper_right; orig_w = float(orig_ur[0]); orig_h = float(orig_ur[1])
                page.cropbox.lower_left = (left, bottom); page.cropbox.upper_right = (orig_w - right, orig_h - top)
                writer = PdfWriter(); writer.add_page(page); temp = io.BytesIO(); writer.write(temp)
                st.markdown("### Cropped Preview (Page 1)")
                thumb = get_page_thumbnail(temp.getvalue(), poppler_path, width=400)
                if thumb: st.image(thumb, width=400)
                file.seek(0)
            except Exception as e: st.error(f"Preview Error: {e}")
            if st.button("Crop & Download", type="primary"):
                reader = PdfReader(file); writer = PdfWriter()
                for page in reader.pages:
                    orig_ur = page.mediabox.upper_right; w, h = float(orig_ur[0]), float(orig_ur[1])
                    page.cropbox.lower_left = (left, bottom); page.cropbox.upper_right = (w - right, h - top); writer.add_page(page)
                out = io.BytesIO(); writer.write(out)
                st.download_button("Download Cropped PDF", out.getvalue(), "cropped.pdf", "application/pdf")

    elif tool == "Sign PDF":
        st.header("✍️ Sign PDF")
        st.write("Create or upload a signature and place it on your document.")
        file = st.file_uploader("Upload PDF", type="pdf")
        if file:
            file_hash = f"{file.name}_{file.size}_sign"
            if 'visual_sign_file_hash' not in st.session_state or st.session_state['visual_sign_file_hash'] != file_hash:
                st.session_state['visual_sign_file_hash'] = file_hash; st.session_state['visual_sign_queue'] = []
                file.seek(0); reader = PdfReader(file)
                for i, page in enumerate(reader.pages):
                    writer = PdfWriter(); writer.add_page(page); p_bytes = io.BytesIO(); writer.write(p_bytes)
                    st.session_state['visual_sign_queue'].append({
                        'id': str(uuid.uuid4()), 'page_num': i + 1, 'bytes': p_bytes.getvalue(), 'rotation': 0
                    })
            with st.expander("👁️ Organize Pages (Rotate / Reorder)", expanded=False):
                if st.session_state['visual_sign_queue']:
                    total_pg = len(st.session_state['visual_sign_queue']); cols = st.columns(4)
                    for i, item in enumerate(st.session_state['visual_sign_queue']):
                        with cols[i % 4]:
                            with st.container():
                                st.markdown(f"<div class='page-card'>", unsafe_allow_html=True)
                                st.caption(f"Pg {i+1}")
                                thumb = get_page_thumbnail(item['bytes'], poppler_path)
                                rot = item.get('rotation', 0)
                                if thumb and rot != 0: thumb = thumb.rotate(-rot, expand=True)
                                if thumb: st.image(thumb, use_container_width=True)
                                c_rot1, c_rot2 = st.columns(2)
                                if c_rot1.button("⟲", key=f"ccw_sg_{item['id']}"): st.session_state['visual_sign_queue'][i]['rotation'] = (rot - 90) % 360; st.rerun()
                                if c_rot2.button("⟳", key=f"cw_sg_{item['id']}"): st.session_state['visual_sign_queue'][i]['rotation'] = (rot + 90) % 360; st.rerun()
                                c1, c3 = st.columns(2)
                                if c1.button("⬅️", key=f"L_sg_{item['id']}") and i > 0: st.session_state['visual_sign_queue'][i], st.session_state['visual_sign_queue'][i-1] = st.session_state['visual_sign_queue'][i-1], st.session_state['visual_sign_queue'][i]; st.rerun()
                                if c3.button("➡️", key=f"R_sg_{item['id']}") and i < total_pg - 1: st.session_state['visual_sign_queue'][i], st.session_state['visual_sign_queue'][i+1] = st.session_state['visual_sign_queue'][i+1], st.session_state['visual_sign_queue'][i]; st.rerun()
                                st.markdown("</div>", unsafe_allow_html=True)
            st.markdown("---")
        
        sig_source = st.radio("Signature Source", ["Draw New", "Type Text", "Upload Image", "Use Default"], horizontal=True)
        final_sig_image = None

        if sig_source == "Draw New":
            if HAS_CANVAS_SUPPORT:
                st.write("Draw your signature below:")
                stroke_width = st.slider("Pencil Thickness", 1, 20, 2)
                canvas_result = st_canvas(stroke_width=stroke_width, stroke_color="black", background_color="rgba(0,0,0,0)", update_streamlit=True, height=150, key="signature_canvas")
                if canvas_result.image_data is not None:
                    final_sig_image = Image.fromarray(canvas_result.image_data.astype('uint8'), 'RGBA')
                    if st.button("Save as Default Signature"): final_sig_image.save("default_signature.png"); st.success("Saved!")
            else: st.warning("Drawing requires `streamlit-drawable-canvas`.")
        elif sig_source == "Type Text":
            sig_text = st.text_input("Enter Signature Text", "John Doe")
            sig_color = st.color_picker("Signature Color", "#000000")
            sig_font_file = st.file_uploader("Upload Font (Optional)", type=["ttf", "otf"])
            if sig_text:
                try:
                    font_size = 100; font = None
                    if sig_font_file: font = ImageFont.truetype(sig_font_file, font_size)
                    else:
                         try: font = ImageFont.truetype("arial.ttf", font_size)
                         except: font = ImageFont.load_default()
                    dummy = Image.new("RGBA", (1, 1)); draw = ImageDraw.Draw(dummy)
                    if hasattr(draw, "textbbox"): bbox = draw.textbbox((0, 0), sig_text, font=font); w = bbox[2] - bbox[0] + 40; h = bbox[3] - bbox[1] + 40
                    else: w, h = draw.textsize(sig_text, font=font); w += 40; h += 40
                    img = Image.new("RGBA", (w, h), (255, 255, 255, 0)); draw = ImageDraw.Draw(img); draw.text((20, 20), sig_text, font=font, fill=sig_color)
                    final_sig_image = img; st.image(final_sig_image, caption="Generated Signature", width=200)
                    if st.button("Save as Default Signature"): final_sig_image.save("default_signature.png"); st.success("Saved!")
                except Exception as e: st.error(f"Error creating signature: {e}")
        elif sig_source == "Upload Image":
            uploaded_sig = st.file_uploader("Upload Signature (PNG/JPG)", type=["png", "jpg", "jpeg"])
            if uploaded_sig:
                final_sig_image = Image.open(uploaded_sig)
                if st.button("Save this as Default"): final_sig_image.save("default_signature.png"); st.success("Saved!")
        elif sig_source == "Use Default":
            if os.path.exists("default_signature.png"): final_sig_image = Image.open("default_signature.png"); st.image(final_sig_image, caption="Default Signature", width=200)
            else: st.error("No default signature found.")

        if file and final_sig_image and st.session_state['visual_sign_queue']:
            st.markdown("### Position & Sign")
            col_preview, col_controls = st.columns([2, 1])
            total_pages = len(st.session_state['visual_sign_queue'])
            with col_controls:
                st.write("**Apply To**")
                sign_scope = st.radio("Pages:", ["All Pages", "Specific Pages"])
                p_input = "1"
                if sign_scope == "Specific Pages": p_input = st.text_input("Page Numbers (e.g. 1, 3-5)", "1")
                if sign_scope == "All Pages": target_indices = list(range(total_pages))
                else: target_indices = parse_order_string(p_input, total_pages)
                st.markdown("---")
                st.write("**Navigation**")
                if target_indices and len(target_indices) > 0:
                    preview_options = [i + 1 for i in target_indices]
                    selected_page_num = st.selectbox("Preview Page", preview_options)
                    preview_page_idx = selected_page_num - 1
                else: st.warning("No valid pages selected."); preview_page_idx = None
                st.markdown("---")
                st.write("**Positioning**")
                x_pos = st.slider("X (Left/Right)", 0, 600, 100)
                y_pos = st.slider("Y (Up/Down)", 0, 800, 100)
                width = st.slider("Width", 50, 400, 150)
                height = st.slider("Height", 20, 200, 50)
                preview_zoom = st.slider("Preview Zoom", 200, 1000, 500, step=50)
            with col_preview:
                if preview_page_idx is not None:
                    try:
                        item = st.session_state['visual_sign_queue'][preview_page_idx]
                        p_reader = PdfReader(io.BytesIO(item['bytes'])); page = p_reader.pages[0]
                        rot = item.get('rotation', 0)
                        if rot != 0: page.rotate(rot)
                        temp_writer = PdfWriter(); temp_writer.add_page(page); temp_pdf_bytes_io = io.BytesIO(); temp_writer.write(temp_pdf_bytes_io); temp_pdf_bytes = temp_pdf_bytes_io.getvalue()
                        temp_reader = PdfReader(io.BytesIO(temp_pdf_bytes)); preview_page = temp_reader.pages[0]; pg_w = float(preview_page.mediabox.width); pg_h = float(preview_page.mediabox.height)
                        packet = io.BytesIO(); c = canvas.Canvas(packet, pagesize=(pg_w, pg_h))
                        img_byte_arr = io.BytesIO(); final_sig_image.save(img_byte_arr, format='PNG'); img_byte_arr.seek(0); sig_img = ImageReader(img_byte_arr)
                        c.drawImage(sig_img, x_pos, y_pos, width=width, height=height, mask='auto')
                        c.save(); packet.seek(0)
                        sig_layer = PdfReader(packet).pages[0]; preview_page.merge_page(sig_layer)
                        final_prev_writer = PdfWriter(); final_prev_writer.add_page(preview_page); prev_bytes = io.BytesIO(); final_prev_writer.write(prev_bytes)
                        thumb = get_page_thumbnail(prev_bytes.getvalue(), poppler_path, width=preview_zoom)
                        if thumb: st.image(thumb, caption=f"Live Preview (Page {preview_page_idx+1})", width=preview_zoom)
                    except Exception as e: st.error(f"Preview Error: {e}")
                else: st.info("Please select pages to view preview.")
            if st.button("Sign & Download PDF", type="primary"):
                try:
                    if sign_scope == "All Pages": target_indices = range(total_pages)
                    else:
                        target_indices = parse_order_string(p_input, total_pages)
                        if not target_indices: st.error("Invalid page selection."); st.stop()
                    writer = PdfWriter()
                    img_byte_arr = io.BytesIO(); final_sig_image.save(img_byte_arr, format='PNG'); img_byte_arr.seek(0)
                    for i, item in enumerate(st.session_state['visual_sign_queue']):
                        p_reader = PdfReader(io.BytesIO(item['bytes'])); page = p_reader.pages[0]
                        rot = item.get('rotation', 0)
                        if rot != 0: page.rotate(rot)
                        if i in target_indices:
                            t_writer = PdfWriter(); t_writer.add_page(page); t_io = io.BytesIO(); t_writer.write(t_io)
                            t_reader = PdfReader(io.BytesIO(t_io.getvalue())); target_page = t_reader.pages[0]
                            pg_w = float(target_page.mediabox.width); pg_h = float(target_page.mediabox.height)
                            packet = io.BytesIO(); c = canvas.Canvas(packet, pagesize=(pg_w, pg_h))
                            img_byte_arr.seek(0); sig_img_rl = ImageReader(img_byte_arr)
                            c.drawImage(sig_img_rl, x_pos, y_pos, width=width, height=height, mask='auto')
                            c.save(); packet.seek(0)
                            sig_layer = PdfReader(packet).pages[0]; target_page.merge_page(sig_layer); writer.add_page(target_page)
                        else: writer.add_page(page)
                    out = io.BytesIO(); writer.write(out)
                    st.download_button("Download Signed PDF", out.getvalue(), "signed_document.pdf", "application/pdf")
                except Exception as e: st.error(f"Error signing document: {e}")

    elif tool == "Lock PDF":
        st.header("🔒 Lock PDF")
        file = st.file_uploader("Upload PDF", type="pdf")
        pw = st.text_input("Password", type="password")
        if file and pw and st.button("Encrypt"):
            reader = PdfReader(file); writer = PdfWriter()
            for page in reader.pages: writer.add_page(page)
            writer.encrypt(pw)
            out = io.BytesIO(); writer.write(out)
            st.download_button("Download", out.getvalue(), "protected.pdf", "application/pdf")

    elif tool == "Decrypt / Unlock PDF":
        st.header("🔓 Decrypt / Unlock PDF")
        decrypt_action = st.radio("Action", ["Unlock with Password", "Recover Lost Password"])
        file = st.file_uploader("Upload Locked PDF", type="pdf")
        if decrypt_action == "Unlock with Password":
            user_pw = st.text_input("Enter Password", type="password")
            if 'unlocked_file_data' not in st.session_state: st.session_state['unlocked_file_data'] = None
            if file and st.button("Unlock"):
                file.seek(0) 
                try:
                    locked_bytes = io.BytesIO(file.read())
                    try:
                        with pikepdf.open(locked_bytes, password=user_pw) as pdf:
                            output_stream = io.BytesIO(); pdf.save(output_stream); st.session_state['unlocked_file_data'] = output_stream.getvalue()
                        st.success("Unlocked successfully!")
                    except pikepdf.PasswordError: st.error("Incorrect password.")
                    except Exception as e: st.error(f"Error: {e}")
                except Exception as e: st.error(f"System Error: {e}")
            if st.session_state['unlocked_file_data']: st.download_button(label="Download Unlocked PDF", data=st.session_state['unlocked_file_data'], file_name="unlocked.pdf", mime="application/pdf")
        else: 
            mode = st.radio("Recovery Mode", ["Dictionary Attack (Fast)", "Brute Force (Comprehensive)"])
            use_custom_list = st.checkbox("Use Custom Wordlist (.txt)")
            custom_words = []
            if use_custom_list:
                wordlist_file = st.file_uploader("Upload Dictionary", type=["txt", "csv", "xlsx", "docx"])
                if wordlist_file: custom_words = read_wordlist_file(wordlist_file); st.info(f"Loaded {len(custom_words)} passwords.")
            charset = ""; max_len = 4
            if mode == "Brute Force (Comprehensive)":
                st.warning("⚠️ Brute force is extremely slow in Python.")
                c1, c2, c3, c4 = st.columns(4)
                use_digits = c1.checkbox("Digits (0-9)", value=True); use_lower = c2.checkbox("Lowercase (a-z)")
                use_upper = c3.checkbox("Uppercase (A-Z)"); use_symbols = c4.checkbox("Symbols")
                max_len = st.slider("Max Length", 1, 8, 4)
                if use_digits: charset += string.digits
                if use_lower: charset += string.ascii_lowercase
                if use_upper: charset += string.ascii_uppercase
                if use_symbols: charset += string.punctuation
            COMMON_PASSWORDS = ["", "123456", "password", "1234", "12345", "12345678", "123456789", "1234567890", "000000", "111111", "password123", "admin", "root", "user", "pdf", "document", "master", "0000", "1234567", "123123"] + [str(y) for y in range(1980, 2030)]
            if file and st.button("Start Recovery"):
                progress_bar = st.progress(0); status_text = st.empty()
                with open("temp_locked.pdf", "wb") as f: f.write(file.read())
                found_password = None; unlocked = False
                def password_generator():
                    yield ""
                    if mode.startswith("Dictionary"):
                        if use_custom_list and custom_words:
                            for pw in custom_words: yield pw
                        for pw in COMMON_PASSWORDS: yield pw
                    else:
                        if not charset: return
                        for length in range(1, max_len + 1):
                            for p in itertools.product(charset, repeat=length): yield "".join(p)
                if mode.startswith("Dictionary"): password_list = list(password_generator()); total_attempts = len(password_list)
                else: password_list = None; total_attempts = sum(len(charset)**l for l in range(1, max_len+1))
                start_time = time.time()
                if password_list:
                    status_text.text(f"Testing {total_attempts} passwords using 4 threads...")
                    batch_size = 50; chunks = [password_list[i:i + batch_size] for i in range(0, len(password_list), batch_size)]
                    with concurrent.futures.ThreadPoolExecutor(max_workers=4) as executor:
                        futures = [executor.submit(check_password_batch, "temp_locked.pdf", chunk) for chunk in chunks]
                        for i, future in enumerate(concurrent.futures.as_completed(futures)):
                            result = future.result()
                            if result is not None: found_password = result if result else "(None - Owner Restriction only)"; unlocked = True; executor.shutdown(wait=False); break
                            progress = min((i + 1) / len(chunks), 1.0); progress_bar.progress(progress)
                else:
                    count = 0
                    for pw in password_generator():
                        count += 1
                        if count % 10 == 0: elapsed = time.time() - start_time; rate = count / elapsed if elapsed > 0 else 0; status_text.text(f"Testing: {pw} | Rate: {rate:.1f} pwd/sec"); progress = min(count / total_attempts, 1.0); progress_bar.progress(progress)
                        try:
                            with pikepdf.open("temp_locked.pdf", password=pw) as pdf: found_password = pw if pw else "(None - Owner Restriction only)"; unlocked = True; break
                        except: continue
                if unlocked:
                    progress_bar.progress(100); status_text.success("Done!")
                    real_pw = "" if found_password == "(None - Owner Restriction only)" else found_password
                    with pikepdf.open("temp_locked.pdf", password=real_pw) as pdf: pdf.save("temp_unlocked.pdf")
                    st.success(f"🔓 Success! PDF Unlocked."); st.info(f"🔑 Password Found: **{found_password}**")
                    with open("temp_unlocked.pdf", "rb") as f: st.download_button("Download Passwordless PDF", f.read(), "unlocked.pdf", "application/pdf")
                    if os.path.exists("temp_unlocked.pdf"): os.remove("temp_unlocked.pdf")
                else: st.error("❌ Failed to find password with current settings.")
                if os.path.exists("temp_locked.pdf"): os.remove("temp_locked.pdf")